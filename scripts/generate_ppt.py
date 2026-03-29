#!/usr/bin/env python3
"""사내 코딩 어시스턴트 도입 방식 제안서 PPT.

기존 제안(벡터 RAG + 웹 UI)에 대한 대안:
구조 기반 MCP + Claude Code
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

OUTPUT = Path(__file__).resolve().parent.parent / "docs" / "proposal.pptx"
OUTPUT.parent.mkdir(parents=True, exist_ok=True)

# 색상
BLACK = RGBColor(0x09, 0x09, 0x0B)
DARK = RGBColor(0x13, 0x13, 0x16)
SURFACE = RGBColor(0x1A, 0x1A, 0x1E)
BORDER = RGBColor(0x25, 0x25, 0x29)
WHITE = RGBColor(0xFA, 0xFA, 0xFA)
GRAY = RGBColor(0xA1, 0xA1, 0xAA)
LIGHT_GRAY = RGBColor(0xD4, 0xD4, 0xD8)
DIM = RGBColor(0x52, 0x52, 0x5B)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
RED = RGBColor(0xEF, 0x44, 0x44)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
CYAN = RGBColor(0x22, 0xD3, 0xEE)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK


def add_text(slide, left, top, width, height, text, size=14, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_rect(slide, left, top, width, height, fill_color=SURFACE, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, items, title_color=AMBER):
    add_rect(slide, left, top, width, height, SURFACE, BORDER)
    add_text(slide, left + 0.2, top + 0.15, width - 0.4, 0.4, title,
             size=13, color=title_color, bold=True)
    y = top + 0.55
    for item in items:
        add_text(slide, left + 0.2, y, width - 0.4, 0.3, item, size=10, color=LIGHT_GRAY)
        y += 0.28


# ================================================================
# 1. 표지
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 1, 1.2, 11, 0.5, "사내 코딩 어시스턴트 도입 방식 검토", size=16, color=GRAY)
add_text(slide, 1, 1.9, 11, 1.2,
         "벡터 RAG + 웹 UI 방식의 한계와\n구조 기반 MCP 방식 제안",
         size=36, color=WHITE, bold=True)

add_text(slide, 1, 3.8, 11, 0.5,
         "Pro*C / ProFrame 코드베이스에 최적화된 접근 방식",
         size=16, color=AMBER)

add_rect(slide, 1, 4.8, 3, 0.5, SURFACE, GREEN)
add_text(slide, 1.15, 4.85, 2.7, 0.4, "PoC 비교 검증 완료", size=12, color=GREEN, bold=True)

add_text(slide, 1, 6.2, 11, 0.4, "2026.03", size=13, color=GRAY)


# ================================================================
# 2. 목차
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 1, 0.8, 11, 0.6, "Contents", size=28, color=WHITE, bold=True)

items = [
    ("01", "현재 제안 검토", "벡터 RAG + 웹 UI 코딩 어시스턴트 방식의 구조"),
    ("02", "Pro*C 환경에서의 한계", "벡터 RAG가 Pro*C 코드베이스에서 왜 작동하지 않는가"),
    ("03", "대안: 구조 기반 MCP", "ProFrame의 구조 정보를 그대로 활용하는 접근"),
    ("04", "PoC 비교 결과", "동일 코드베이스에서 두 방식 정량 비교"),
    ("05", "구조 기반 방식의 실제 동작", "개발자 시나리오별 시연"),
    ("06", "도입 방안 비교", "구현 복잡도, 비용, 유지보수 비교"),
    ("07", "제안", "권장 접근 방식 및 다음 단계"),
]

for i, (num, title, desc) in enumerate(items):
    y = 1.8 + i * 0.7
    add_text(slide, 1.5, y, 0.6, 0.4, num, size=20, color=AMBER, bold=True,
             font_name="JetBrains Mono")
    add_text(slide, 2.3, y, 4, 0.35, title, size=16, color=WHITE, bold=True)
    add_text(slide, 2.3, y + 0.3, 8, 0.3, desc, size=11, color=GRAY)


# ================================================================
# 3. 현재 제안 검토
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 1, 0.5, 4, 0.4, "01", size=14, color=AMBER, bold=True)
add_text(slide, 1, 0.9, 11, 0.6, "현재 제안 검토", size=28, color=WHITE, bold=True)
add_text(slide, 1, 1.5, 11, 0.35,
         "벡터 RAG + 별도 웹 UI 코딩 어시스턴트를 구축하는 방식",
         size=13, color=GRAY)

add_card(slide, 1, 2.1, 5.3, 3.2, "제안된 구조", [
    "1. 코드를 텍스트 청크로 분할 (6~50줄)",
    "2. 임베딩 모델로 벡터화 → 벡터 DB 저장",
    "3. 사용자 질문 → 유사도 검색 → 관련 청크 반환",
    "4. LLM이 청크 기반으로 답변 생성",
    "5. 별도 웹 UI에서 사용",
    "",
    "필요 인프라:",
    "  • 임베딩 모델 (GPU 또는 API)",
    "  • 벡터 DB (ChromaDB / Pinecone 등)",
    "  • 웹 서버 + 프론트엔드",
    "  • LLM API 연동",
])

add_card(slide, 7, 2.1, 5.3, 3.2, "이 방식의 전제 조건", [
    "코드가 자연어 텍스트처럼 작동한다고 가정",
    "→ 유사한 단어 = 유사한 기능",
    "",
    "이 전제가 맞는 경우:",
    "  • 일반 텍스트 문서 검색",
    "  • 다양한 언어/프레임워크의 오픈소스 코드",
    "  • 함수명/변수명이 자기 설명적인 코드",
    "",
    "이 전제가 틀리는 경우:",
    "  • 같은 변수명이 다른 도메인에서 반복",
    "  • 비즈니스 로직이 SQL에 있는 코드",
    "  • 호출 관계가 코드 텍스트에 드러나지 않는 코드",
], title_color=RED)

add_rect(slide, 1, 5.7, 11.3, 1, SURFACE, RED)
add_text(slide, 1.2, 5.8, 10.9, 0.4,
         "질문: Pro*C/ProFrame 코드베이스가 이 전제에 맞는가?",
         size=15, color=RED, bold=True)
add_text(slide, 1.2, 6.2, 10.9, 0.4,
         "→ 다음 장에서 검증합니다.",
         size=12, color=GRAY)


# ================================================================
# 4. Pro*C 환경에서의 한계
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 1, 0.5, 4, 0.4, "02", size=14, color=AMBER, bold=True)
add_text(slide, 1, 0.9, 11, 0.6, "Pro*C 환경에서 벡터 RAG의 한계", size=28, color=WHITE, bold=True)

# 문제 1
add_rect(slide, 1, 1.8, 3.5, 2.2, SURFACE, RED)
add_text(slide, 1.2, 1.9, 3.1, 0.35, "문제 1: 노이즈", size=14, color=RED, bold=True)
add_text(slide, 1.2, 2.35, 3.1, 1.5,
         "tax_amt, acct_no, fee_rate,\nHOLD_QTY, EVAL_AMT ...\n\n"
         "이 변수명들이 주식, 채권, 펀드,\n배당, 대출 등 모든 도메인에서\n동일하게 사용됨.\n\n"
         "→ 벡터 유사도가 높지만\n   실제로는 무관한 코드",
         size=10, color=LIGHT_GRAY)

# 문제 2
add_rect(slide, 4.8, 1.8, 3.5, 2.2, SURFACE, RED)
add_text(slide, 5, 1.9, 3.1, 0.35, "문제 2: 컨텍스트 단절", size=14, color=RED, bold=True)
add_text(slide, 5, 2.35, 3.1, 1.5,
         "Pro*C 함수를 이해하려면:\n"
         "  • 호출하는 다른 함수\n"
         "  • 입출력 구조체 (struct)\n"
         "  • 참조 테이블 DDL\n\n"
         "이 모두 다른 파일에 있음.\n"
         "청크 검색은 한 파일의 조각만\n반환하므로 전체 그림이 안 보임.",
         size=10, color=LIGHT_GRAY)

# 문제 3
add_rect(slide, 8.6, 1.8, 3.7, 2.2, SURFACE, RED)
add_text(slide, 8.8, 1.9, 3.3, 0.35, "문제 3: 구조 낭비", size=14, color=RED, bold=True)
add_text(slide, 8.8, 2.35, 3.3, 1.5,
         "ProFrame이 이미 제공하는 것:\n"
         "  • 서비스 등록 정보 (이름, 설명)\n"
         "  • FDL (입출력 구조체 정의)\n"
         "  • EXECUTE_SERVICE (호출 관계)\n\n"
         "벡터 RAG는 이 구조 정보를\n전부 무시하고 코드를 텍스트\n조각으로 취급함.",
         size=10, color=LIGHT_GRAY)

# 하단 요약
add_rect(slide, 1, 4.4, 11.3, 0.8, SURFACE, AMBER)
add_text(slide, 1.2, 4.45, 10.9, 0.35,
         "Pro*C/ProFrame 코드베이스의 특성",
         size=13, color=AMBER, bold=True)
add_text(slide, 1.2, 4.8, 10.9, 0.3,
         "C + SQL 혼합  |  동일 변수명이 도메인 간 반복  |  비즈니스 로직 = EXEC SQL  |  ProFrame이 이미 구조화",
         size=11, color=LIGHT_GRAY)

add_text(slide, 1, 5.6, 11.3, 0.5,
         "결론: 구조화된 코드에 비정형 텍스트 검색을 적용하는 것은\n"
         "엑셀 테이블이 있는데 자연어로 검색하는 것과 같다.",
         size=14, color=WHITE, bold=True)


# ================================================================
# 5. 대안: 구조 기반 MCP
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 1, 0.5, 4, 0.4, "03", size=14, color=AMBER, bold=True)
add_text(slide, 1, 0.9, 11, 0.6, "대안: 구조 기반 MCP + 기존 AI 코딩 도구 활용", size=26, color=WHITE, bold=True)

# 아키텍처
add_rect(slide, 4.2, 1.8, 5, 0.7, AMBER)
add_text(slide, 4.2, 1.85, 5, 0.6, "Claude Code (기존 AI 코딩 도구)",
         size=14, color=BLACK, bold=True, align=PP_ALIGN.CENTER)

# 3 pillars
add_rect(slide, 1, 3, 3.5, 1.6, SURFACE, BLUE)
add_text(slide, 1.2, 3.1, 3.1, 0.3, "svn-bridge MCP", size=13, color=BLUE, bold=True)
add_text(slide, 1.2, 3.5, 3.1, 1, "SVN 이력 조회\n커밋 로그, diff, blame\n코드 읽기 (svn_cat)", size=10, color=GRAY)

add_rect(slide, 5, 3, 3.5, 1.6, SURFACE, AMBER)
add_text(slide, 5.2, 3.1, 3.1, 0.3, "proc-retrieval MCP", size=13, color=AMBER, bold=True)
add_text(slide, 5.2, 3.5, 3.1, 1, "코드 구조 검색 (10개 Tool)\n호출 그래프, 영향도, 업무 흐름\nSQLite 경량 DB", size=10, color=GRAY)

add_rect(slide, 9, 3, 3.3, 1.6, SURFACE, GREEN)
add_text(slide, 9.2, 3.1, 2.9, 0.3, "Read / Edit (내장)", size=13, color=GREEN, bold=True)
add_text(slide, 9.2, 3.5, 2.9, 1, "실제 코드 파일 읽기/수정\nClaude Code 내장 기능\nRepo 원본 직접 접근", size=10, color=GRAY)

# 핵심 차이
add_rect(slide, 1, 5, 5.5, 1.8, SURFACE, BORDER)
add_text(slide, 1.2, 5.1, 5.1, 0.3, "핵심 차이점", size=14, color=AMBER, bold=True)
diffs = [
    "별도 웹 UI 개발 불필요 → Claude Code가 UI",
    "임베딩 모델/GPU 불필요 → SQLite 파일 1개",
    "코드를 청크로 자르지 않음 → 구조 그대로 활용",
    "MCP는 경로만 반환 → 코드는 repo에서 직접 Read",
    "ProFrame 구조 정보를 100% 활용",
]
for i, d in enumerate(diffs):
    add_text(slide, 1.3, 5.5 + i * 0.24, 5, 0.24, f"• {d}", size=10, color=LIGHT_GRAY)

add_rect(slide, 7, 5, 5.3, 1.8, SURFACE, BORDER)
add_text(slide, 7.2, 5.1, 5, 0.3, "MCP(Model Context Protocol)란?", size=14, color=BLUE, bold=True)
add_text(slide, 7.2, 5.5, 5, 1.2,
         "AI 코딩 도구가 외부 데이터 소스에\n"
         "접근할 수 있게 해주는 표준 프로토콜.\n\n"
         "Claude Code, Cursor, Windsurf 등\n"
         "주요 AI 코딩 도구가 지원.\n\n"
         "→ 한 번 만들면 여러 도구에서 사용 가능",
         size=10, color=LIGHT_GRAY)


# ================================================================
# 6. PoC 비교 결과
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 1, 0.5, 4, 0.4, "04", size=14, color=AMBER, bold=True)
add_text(slide, 1, 0.9, 11, 0.6, "PoC 비교 결과", size=28, color=WHITE, bold=True)
add_text(slide, 1, 1.5, 11, 0.35,
         "동일 코드베이스(128파일, 81함수)에서 10개 시나리오로 정량 비교",
         size=13, color=GRAY)

# 비교 테이블
headers = [("지표", 2.5, GRAY), ("벡터 RAG", 2.5, RED), ("구조 기반", 2.5, GREEN), ("의미", 4, GRAY)]
x = 1
for label, w, c in headers:
    add_rect(slide, x, 2.1, w, 0.5, SURFACE, BORDER)
    add_text(slide, x + 0.1, 2.15, w - 0.2, 0.4, label, size=12, color=c, bold=True)
    x += w

rows = [
    ("Recall\n(정답을 찾았는가)", "40.0%", "98.0%", "구조 기반이 2.5배 높은 정답률"),
    ("Noise\n(잘못된 결과 비율)", "22.0%", "0.6%", "벡터는 노이즈가 5건 중 1건"),
    ("Completeness\n(완전한 컨텍스트)", "52.5%", "97.5%", "구조 기반은 호출함수+DDL까지 포함"),
    ("Precision\n(반환 중 정답 비율)", "34.0%", "40.7%", "구조 기반이 유틸 함수도 포함해서 비슷"),
]

for i, (label, rag, struct, meaning) in enumerate(rows):
    y = 2.6 + i * 0.7
    x = 1
    for val, w, c in [(label, 2.5, LIGHT_GRAY), (rag, 2.5, RED), (struct, 2.5, GREEN), (meaning, 4, GRAY)]:
        add_rect(slide, x, y, w, 0.7, DARK, BORDER)
        add_text(slide, x + 0.1, y + 0.1, w - 0.2, 0.5, val, size=10, color=c,
                 bold=(c == GREEN))
        x += w

# 하단 인사이트
add_rect(slide, 1, 5.5, 11.3, 1.3, SURFACE, AMBER)
add_text(slide, 1.2, 5.55, 10.9, 0.35, "왜 이런 차이가 나는가?", size=14, color=AMBER, bold=True)

add_text(slide, 1.2, 6, 5, 0.7,
         "벡터 RAG:\n"
         "\"환매 수수료 계산\" 검색 → tax_amt, fee_rate 등\n"
         "유사 변수명을 가진 채권/배당/대출 코드가 노이즈로 반환",
         size=10, color=LIGHT_GRAY)

add_text(slide, 6.5, 6, 5.5, 0.7,
         "구조 기반:\n"
         "\"환매 수수료 계산\" → RD_FEE_03 특정\n"
         "→ 호출 그래프(CM_TAX_01, CM_DATE_CHK) + DDL 자동 확장",
         size=10, color=LIGHT_GRAY)


# ================================================================
# 7. 실제 동작 시연
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 1, 0.5, 4, 0.4, "05", size=14, color=AMBER, bold=True)
add_text(slide, 1, 0.9, 11, 0.6, "구조 기반 방식의 실제 동작", size=28, color=WHITE, bold=True)

scenarios = [
    ("이 함수 수정하면 영향 어디?",
     "get_impact(\"CM_TAX_01\")",
     "9개 직접 호출자 + 9개 간접 호출자\ndepth별 영향 범위 즉시 파악"),
    ("이 테이블 누가 쓰지?",
     "find_by_table(\"TB_ACCT_MST\")",
     "14개 컬럼, SELECT/INSERT/UPDATE별\n사용 함수 목록"),
    ("이 업무 전체 흐름 보여줘",
     "get_business_flow(\"FD_PAY_01\")",
     "호출 트리 + 각 단계 SQL 연산\n(SELECT TB_FUND_DIST → FD_TAX_02 → ...)"),
    ("이 에러 어디서 나는 거야?",
     "find_by_error_code(\"E101\")",
     "ST_BUY_01 \"주문가능금액 부족\"\n에러 설정 위치 + 주변 코드"),
    ("이 함수 코드 수정해줘",
     "get_function → Read → Edit",
     "구조 파악(MCP) → 코드 로드(Read)\n→ 수정(Edit) — 항상 최신 코드"),
]

for i, (question, tool, result) in enumerate(scenarios):
    y = 1.7 + i * 1.05
    add_rect(slide, 1, y, 11.3, 0.9, SURFACE, BORDER)
    add_text(slide, 1.2, y + 0.05, 3.5, 0.3,
             f"\"{question}\"", size=12, color=WHITE, bold=True)
    add_text(slide, 4.8, y + 0.05, 4, 0.3,
             tool, size=10, color=AMBER, font_name="JetBrains Mono")
    add_text(slide, 4.8, y + 0.35, 7, 0.5,
             f"→ {result}", size=10, color=GREEN)


# ================================================================
# 8. 도입 방안 비교
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 1, 0.5, 4, 0.4, "06", size=14, color=AMBER, bold=True)
add_text(slide, 1, 0.9, 11, 0.6, "도입 방안 비교", size=28, color=WHITE, bold=True)

# 비교 항목
items_compare = [
    ("UI 개발", "웹 프론트엔드 + 백엔드\n별도 구축 필요", "불필요\nClaude Code가 UI", RED, GREEN),
    ("인프라", "임베딩 모델 (GPU/API)\n벡터 DB 서버\n웹 서버", "SQLite 파일 1개\n(수 MB)", RED, GREEN),
    ("코드 최신성", "임베딩 재생성 필요\n(코드 변경 시마다)", "repo에서 직접 Read\n항상 최신", RED, GREEN),
    ("검색 정확도", "Recall 40%\nNoise 22%", "Recall 98%\nNoise 0.6%", RED, GREEN),
    ("유지보수", "임베딩 모델 업데이트\n벡터 DB 관리\n웹 서버 운영", "SQLite 재인덱싱\n(수 초)", RED, GREEN),
    ("구현 기간", "프론트 + 백엔드 + RAG\n파이프라인 = 수 개월", "MCP 2개 + 파서\n= 4주", RED, GREEN),
]

# 헤더
add_rect(slide, 1, 1.7, 2.5, 0.45, SURFACE, BORDER)
add_text(slide, 1.1, 1.75, 2.3, 0.35, "항목", size=11, color=GRAY, bold=True)
add_rect(slide, 3.5, 1.7, 4, 0.45, SURFACE, RED)
add_text(slide, 3.6, 1.75, 3.8, 0.35, "벡터 RAG + 웹 UI", size=11, color=RED, bold=True)
add_rect(slide, 7.5, 1.7, 4.8, 0.45, SURFACE, GREEN)
add_text(slide, 7.6, 1.75, 4.6, 0.35, "구조 기반 MCP + Claude Code", size=11, color=GREEN, bold=True)

for i, (label, rag_val, struct_val, rag_c, struct_c) in enumerate(items_compare):
    y = 2.15 + i * 0.8
    add_rect(slide, 1, y, 2.5, 0.8, DARK, BORDER)
    add_text(slide, 1.1, y + 0.1, 2.3, 0.6, label, size=11, color=LIGHT_GRAY, bold=True)
    add_rect(slide, 3.5, y, 4, 0.8, DARK, BORDER)
    add_text(slide, 3.6, y + 0.1, 3.8, 0.6, rag_val, size=10, color=rag_c)
    add_rect(slide, 7.5, y, 4.8, 0.8, DARK, BORDER)
    add_text(slide, 7.6, y + 0.1, 4.6, 0.6, struct_val, size=10, color=struct_c, bold=True)


# ================================================================
# 9. 제안
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 1, 0.5, 4, 0.4, "07", size=14, color=AMBER, bold=True)
add_text(slide, 1, 0.9, 11, 0.6, "제안", size=28, color=WHITE, bold=True)

# 제안 내용
add_rect(slide, 1, 1.8, 11.3, 1.2, SURFACE, AMBER)
add_text(slide, 1.2, 1.85, 10.9, 0.35,
         "사내 코딩 어시스턴트 도입 시, 벡터 RAG + 별도 웹 UI 구축 대신",
         size=14, color=LIGHT_GRAY)
add_text(slide, 1.2, 2.25, 10.9, 0.5,
         "ProFrame 구조 정보 기반 MCP 서버를 구축하고, 기존 AI 코딩 도구(Claude Code 등)와 연동",
         size=16, color=AMBER, bold=True)

# 필요 리소스
add_card(slide, 1, 3.3, 5.3, 2.5, "필요 리소스", [
    "개발자 1명 (파서 + MCP 구축, 3주)",
    "업무 시니어 1명 (서비스 설명 검토, 2~3일)",
    "",
    "인프라:",
    "  SQLite 파일 1개 (수 MB)",
    "  SVN 접근 권한",
    "  Claude Code 라이선스",
    "",
    "불필요: GPU, 벡터 DB, 웹 서버, 프론트엔드",
], title_color=BLUE)

# 로드맵
add_rect(slide, 7, 3.3, 5.3, 2.5, SURFACE, BORDER)
add_text(slide, 7.2, 3.4, 5, 0.35, "구현 일정 (4주)", size=14, color=GREEN, bold=True)

phases = [
    ("1주", "Pro*C 파서 개발\n실제 .pc 파일 파싱 검증"),
    ("2주", "MCP 서버 구축\nproc-retrieval + svn-bridge"),
    ("3주", "전체 코드베이스 인덱싱\n검색 정확도 튜닝"),
    ("4주", "Claude Code 연동\n개발자 파일럿 테스트"),
]
for i, (week, desc) in enumerate(phases):
    y = 3.9 + i * 0.45
    add_text(slide, 7.3, y, 0.8, 0.35, week, size=11, color=AMBER, bold=True,
             font_name="JetBrains Mono")
    add_text(slide, 8.2, y, 4, 0.4, desc, size=10, color=LIGHT_GRAY)

# 다음 단계
add_rect(slide, 1, 6.1, 11.3, 0.8, SURFACE, GREEN)
add_text(slide, 1.2, 6.15, 10.9, 0.3, "다음 단계", size=13, color=GREEN, bold=True)
add_text(slide, 1.2, 6.5, 10.9, 0.35,
         "1. 실제 Pro*C 파일 10개로 파서 검증    "
         "2. ProFrame 서비스 설명 품질 확인    "
         "3. 파일럿 팀 선정",
         size=12, color=LIGHT_GRAY)


# ================================================================
# 10. 마무리
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 1, 2, 11.3, 0.6,
         "코딩 어시스턴트를 새로 만들 필요가 없습니다.",
         size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, 1, 3, 11.3, 0.6,
         "AI 코딩 도구가 우리 코드를 이해할 수 있게\n구조 정보를 제공하면 됩니다.",
         size=20, color=AMBER, align=PP_ALIGN.CENTER)

add_rect(slide, 3.5, 4.3, 6, 1.8, SURFACE, BORDER)
add_text(slide, 3.7, 4.4, 5.6, 0.3, "PoC 검증 결과", size=13, color=GRAY, bold=True)
add_text(slide, 3.7, 4.8, 5.6, 0.35,
         "Recall 98%  |  Noise 0.6%  |  50/50 테스트 통과",
         size=14, color=GREEN, bold=True)
add_text(slide, 3.7, 5.3, 5.6, 0.3,
         "벡터 RAG 대비: Recall 2.5배 ↑ | Noise 36배 ↓",
         size=12, color=LIGHT_GRAY)
add_text(slide, 3.7, 5.7, 5.6, 0.3,
         "2명 × 4주  |  웹 UI 불필요  |  GPU/벡터DB 불필요",
         size=12, color=LIGHT_GRAY)


# ================================================================
# 저장
# ================================================================
prs.save(str(OUTPUT))
print(f"PPT 생성 완료: {OUTPUT}")


if __name__ == "__main__":
    main() if "main" in dir() else None
